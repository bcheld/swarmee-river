from __future__ import annotations

import json
import re
from pathlib import Path
from typing import Any

from strands import tool

from swarmee_river.tool_permissions import set_permissions
from swarmee_river.utils.path_utils import safe_cwd
from swarmee_river.utils.text_utils import truncate

_SUPPORTED_EXTENSIONS = {".docx", ".xlsx", ".pptx"}


def _success(text: str, *, max_chars: int) -> dict[str, Any]:
    return {"status": "success", "content": [{"text": truncate(text, max_chars)}]}


def _error(text: str, *, max_chars: int) -> dict[str, Any]:
    return {"status": "error", "content": [{"text": truncate(text, max_chars)}]}


def _resolve_path(path: str, *, cwd: str | None, for_write: bool) -> tuple[Path, Path]:
    rel_path = (path or "").strip()
    if not rel_path:
        raise ValueError("path is required")

    base = safe_cwd(cwd)
    target = (base / rel_path).expanduser().resolve(strict=False)
    if base not in target.parents and target != base:
        raise ValueError("Refusing to operate outside cwd")

    if not for_write:
        if not target.exists() or not target.is_file():
            raise FileNotFoundError(f"File not found: {rel_path}")
        target = target.resolve()
        if base not in target.parents and target != base:
            raise ValueError("Refusing to operate outside cwd")

    return base, target


def _read_docx(path: Path, *, max_chars: int) -> dict[str, Any]:
    try:
        from docx import Document
        from docx.document import Document as DocxDocument
        from docx.table import Table, _Cell
        from docx.text.paragraph import Paragraph
    except ImportError:
        return _error(
            "python-docx is required. Install with: pip install python-docx",
            max_chars=max_chars,
        )

    def _iter_blocks(parent: Any) -> Any:
        if isinstance(parent, DocxDocument):
            parent_element = parent.element.body
        elif isinstance(parent, _Cell):
            parent_element = parent._tc
        else:
            return
        for child in parent_element.iterchildren():
            if child.tag.endswith("}p"):
                yield Paragraph(child, parent)
            elif child.tag.endswith("}tbl"):
                yield Table(child, parent)

    def _table_to_markdown(table: Any) -> str:
        rows: list[list[str]] = []
        for row in table.rows:
            cells = [str(cell.text or "").replace("\n", " ").strip() for cell in row.cells]
            rows.append(cells)
        if not rows:
            return ""

        width = max(len(r) for r in rows)
        normalized: list[list[str]] = []
        for row in rows:
            padded = row + [""] * (width - len(row))
            normalized.append([c.replace("|", "\\|") for c in padded])

        header = normalized[0]
        body = normalized[1:] if len(normalized) > 1 else []
        lines = [
            "| " + " | ".join(header) + " |",
            "| " + " | ".join(["---"] * width) + " |",
        ]
        for row in body:
            lines.append("| " + " | ".join(row) + " |")
        return "\n".join(lines)

    try:
        doc = Document(str(path))
    except Exception as e:
        return _error(f"Failed to read .docx file: {e}", max_chars=max_chars)

    blocks: list[str] = []
    for block in _iter_blocks(doc):
        if isinstance(block, Paragraph):
            text = str(block.text or "").strip()
            style_name = str(getattr(getattr(block, "style", None), "name", "") or "").strip().lower()

            image_lines: list[str] = []
            try:
                for node in block._element.xpath(".//*[local-name()='docPr']"):
                    desc = (
                        str(node.get("descr") or node.get("title") or node.get("name") or "").strip()  # noqa: PLC2801
                    )
                    if desc:
                        image_lines.append(f"[Image: {desc}]")
                    else:
                        image_lines.append("[Image]")
            except Exception:
                pass

            if style_name.startswith("heading") and text:
                match = re.search(r"(\d+)", style_name)
                level = int(match.group(1)) if match else 1
                level = min(6, max(1, level))
                blocks.append(f"{'#' * level} {text}")
            elif "list bullet" in style_name and text:
                blocks.append(f"- {text}")
            elif "list number" in style_name and text:
                blocks.append(f"1. {text}")
            elif text:
                blocks.append(text)

            blocks.extend(image_lines)
        elif isinstance(block, Table):
            table_md = _table_to_markdown(block)
            if table_md:
                blocks.append(table_md)

    rendered = "\n\n".join(b for b in blocks if b.strip())
    if not rendered.strip():
        rendered = "(no readable text content)"
    return _success(rendered, max_chars=max_chars)


def _read_xlsx(path: Path, *, sheet_name: str | None, max_rows: int, max_chars: int) -> dict[str, Any]:
    try:
        from openpyxl import load_workbook
    except ImportError:
        return _error(
            "openpyxl is required. Install with: pip install openpyxl",
            max_chars=max_chars,
        )

    capped_rows = max(1, min(1000, int(max_rows)))

    def _cell_to_text(value: Any) -> str:
        if value is None:
            return ""
        if isinstance(value, float):
            return f"{value:.2f}"
        return str(value)

    try:
        workbook = load_workbook(filename=str(path), data_only=True, read_only=True)
    except Exception as e:
        return _error(f"Failed to read .xlsx file: {e}", max_chars=max_chars)

    try:
        if sheet_name:
            if sheet_name not in workbook.sheetnames:
                return _error(
                    f"Sheet not found: {sheet_name}. Available: {', '.join(workbook.sheetnames)}",
                    max_chars=max_chars,
                )
            ws = workbook[sheet_name]
        else:
            ws = workbook.active

        header_values = next(ws.iter_rows(min_row=1, max_row=1, values_only=True), None)
        if header_values is None:
            return _success(f"# Sheet: {ws.title}\n\n(no rows)", max_chars=max_chars)

        headers = [(_cell_to_text(v) or f"Column {idx + 1}") for idx, v in enumerate(header_values)]
        width = len(headers)
        if width == 0:
            return _success(f"# Sheet: {ws.title}\n\n(no columns)", max_chars=max_chars)

        data_rows: list[list[str]] = []
        for row in ws.iter_rows(min_row=2, max_row=1 + capped_rows, values_only=True):
            rendered_row = [_cell_to_text(v).replace("|", "\\|") for v in row]
            if len(rendered_row) < width:
                rendered_row.extend([""] * (width - len(rendered_row)))
            elif len(rendered_row) > width:
                rendered_row = rendered_row[:width]
            data_rows.append(rendered_row)

        lines: list[str] = [
            f"# Sheet: {ws.title}",
            "",
            "| " + " | ".join(h.replace("|", "\\|") for h in headers) + " |",
            "| " + " | ".join(["---"] * width) + " |",
        ]
        for row in data_rows:
            lines.append("| " + " | ".join(row) + " |")

        total_data_rows = max(0, int(ws.max_row or 0) - 1)
        if total_data_rows > len(data_rows):
            lines.append("")
            lines.append(f"... ({total_data_rows - len(data_rows)} more rows)")

        return _success("\n".join(lines).strip(), max_chars=max_chars)
    finally:
        try:
            workbook.close()
        except Exception:
            pass


def _read_pptx(path: Path, *, max_chars: int) -> dict[str, Any]:
    try:
        from pptx import Presentation
    except ImportError:
        return _error(
            "python-pptx is required. Install with: pip install python-pptx",
            max_chars=max_chars,
        )

    try:
        presentation = Presentation(str(path))
    except Exception as e:
        return _error(f"Failed to read .pptx file: {e}", max_chars=max_chars)

    blocks: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        title = ""
        try:
            if slide.shapes.title is not None and slide.shapes.title.text:
                title = slide.shapes.title.text.strip()
        except Exception:
            title = ""
        if not title:
            title = "Untitled"

        body_parts: list[str] = []
        for shape in slide.shapes:
            try:
                if not getattr(shape, "has_text_frame", False):
                    continue
                text = str(shape.text or "").strip()
                if not text:
                    continue
                if slide.shapes.title is not None and shape == slide.shapes.title:
                    continue
                body_parts.append(text)
            except Exception:
                continue

        lines = [f"## Slide {index}: {title}"]
        if body_parts:
            lines.append("\n".join(body_parts))

        notes_text = ""
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
                notes_text = str(slide.notes_slide.notes_text_frame.text or "").strip()
        except Exception:
            notes_text = ""
        if notes_text:
            lines.append(f"> Note: {notes_text}")

        blocks.append("\n".join(lines).strip())

    rendered = "\n\n".join(blocks).strip() if blocks else "(no slides)"
    return _success(rendered, max_chars=max_chars)


def _write_docx(path: Path, *, content: Any, max_chars: int) -> dict[str, Any]:
    try:
        from docx import Document
    except ImportError:
        return _error(
            "python-docx is required. Install with: pip install python-docx",
            max_chars=max_chars,
        )

    if content is None:
        return _error("content is required for action=write", max_chars=max_chars)
    markdown = str(content)

    def _parse_table(lines: list[str], start: int) -> tuple[int, list[list[str]]]:
        rows: list[list[str]] = []
        i = start
        while i < len(lines):
            raw = lines[i].strip()
            if not raw.startswith("|"):
                break
            body = raw.strip("|")
            cells = [c.strip().replace("\\|", "|") for c in body.split("|")]
            rows.append(cells)
            i += 1
        if len(rows) >= 2:
            sep = rows[1]
            if all(re.fullmatch(r"[:\- ]+", cell) for cell in sep):
                rows.pop(1)
        return i, rows

    try:
        doc = Document()
        lines = markdown.splitlines()
        idx = 0
        while idx < len(lines):
            line = lines[idx].rstrip()
            stripped = line.strip()
            if not stripped:
                idx += 1
                continue

            heading_match = re.match(r"^(#{1,6})\s+(.*)$", stripped)
            if heading_match:
                level = len(heading_match.group(1))
                text = heading_match.group(2).strip()
                doc.add_heading(text, level=level)
                idx += 1
                continue

            if stripped.startswith("|"):
                next_idx, table_rows = _parse_table(lines, idx)
                if table_rows:
                    max_cols = max(len(r) for r in table_rows)
                    table = doc.add_table(rows=len(table_rows), cols=max_cols)
                    for r_idx, row in enumerate(table_rows):
                        padded = row + [""] * (max_cols - len(row))
                        for c_idx, value in enumerate(padded):
                            table.cell(r_idx, c_idx).text = value
                    idx = next_idx
                    continue

            bullet_match = re.match(r"^[-*]\s+(.*)$", stripped)
            if bullet_match:
                doc.add_paragraph(bullet_match.group(1).strip(), style="List Bullet")
                idx += 1
                continue

            number_match = re.match(r"^\d+\.\s+(.*)$", stripped)
            if number_match:
                doc.add_paragraph(number_match.group(1).strip(), style="List Number")
                idx += 1
                continue

            doc.add_paragraph(stripped)
            idx += 1

        path.parent.mkdir(parents=True, exist_ok=True)
        doc.save(str(path))
    except Exception as e:
        return _error(f"Failed to write .docx file: {e}", max_chars=max_chars)

    size = path.stat().st_size if path.exists() else 0
    return _success(f"Wrote {size} bytes to {path}", max_chars=max_chars)


def _write_xlsx(path: Path, *, content: Any, sheet_name: str | None, max_chars: int) -> dict[str, Any]:
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font
        from openpyxl.utils import get_column_letter
    except ImportError:
        return _error(
            "openpyxl is required. Install with: pip install openpyxl",
            max_chars=max_chars,
        )

    if content is None:
        return _error("content is required for action=write", max_chars=max_chars)

    try:
        rows_data = json.loads(content) if isinstance(content, str) else content
    except Exception as e:
        return _error(f"content must be valid JSON for .xlsx write: {e}", max_chars=max_chars)

    if not isinstance(rows_data, list):
        return _error("content must be a JSON array of arrays for .xlsx write", max_chars=max_chars)

    try:
        wb = Workbook()
        ws = wb.active
        ws.title = (sheet_name or "Sheet1").strip() or "Sheet1"

        for r_idx, row in enumerate(rows_data, start=1):
            if isinstance(row, (list, tuple)):
                values = list(row)
            else:
                values = [row]
            for c_idx, value in enumerate(values, start=1):
                ws.cell(row=r_idx, column=c_idx, value=value)

        if rows_data:
            first_row = rows_data[0]
            header_width = len(first_row) if isinstance(first_row, (list, tuple)) else 1
            for c_idx in range(1, header_width + 1):
                ws.cell(row=1, column=c_idx).font = Font(bold=True)

        for c_idx in range(1, ws.max_column + 1):
            max_len = 0
            for r_idx in range(1, ws.max_row + 1):
                value = ws.cell(row=r_idx, column=c_idx).value
                length = len(str(value)) if value is not None else 0
                if length > max_len:
                    max_len = length
            ws.column_dimensions[get_column_letter(c_idx)].width = min(max(8, max_len + 2), 60)

        path.parent.mkdir(parents=True, exist_ok=True)
        wb.save(str(path))
    except Exception as e:
        return _error(f"Failed to write .xlsx file: {e}", max_chars=max_chars)

    size = path.stat().st_size if path.exists() else 0
    return _success(f"Wrote {size} bytes to {path}", max_chars=max_chars)


def _write_pptx(path: Path, *, content: Any, max_chars: int) -> dict[str, Any]:
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        return _error(
            "python-pptx is required. Install with: pip install python-pptx",
            max_chars=max_chars,
        )

    if content is None:
        return _error("content is required for action=write", max_chars=max_chars)

    try:
        slides_data = json.loads(content) if isinstance(content, str) else content
    except Exception as e:
        return _error(f"content must be valid JSON for .pptx write: {e}", max_chars=max_chars)

    if not isinstance(slides_data, list):
        return _error("content must be a JSON array of {title, body, notes?} for .pptx write", max_chars=max_chars)

    try:
        prs = Presentation()
        layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]

        for raw_slide in slides_data:
            if not isinstance(raw_slide, dict):
                return _error("each slide item must be an object with title/body/notes fields", max_chars=max_chars)
            title = str(raw_slide.get("title", "") or "")
            body = str(raw_slide.get("body", "") or "")
            notes = str(raw_slide.get("notes", "") or "").strip()

            slide = prs.slides.add_slide(layout)

            if slide.shapes.title is not None:
                slide.shapes.title.text = title
            else:
                title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(8.4), Inches(0.8))
                title_box.text_frame.text = title

            body_frame = None
            for shape in slide.placeholders:
                if slide.shapes.title is not None and shape == slide.shapes.title:
                    continue
                if getattr(shape, "has_text_frame", False):
                    body_frame = shape.text_frame
                    break

            if body_frame is None:
                body_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(4.5))
                body_frame = body_box.text_frame

            body_frame.clear()
            body_frame.text = body

            if notes:
                notes_frame = slide.notes_slide.notes_text_frame
                notes_frame.clear()
                notes_frame.text = notes

        path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(path))
    except Exception as e:
        return _error(f"Failed to write .pptx file: {e}", max_chars=max_chars)

    size = path.stat().st_size if path.exists() else 0
    return _success(f"Wrote {size} bytes to {path}", max_chars=max_chars)


def _info_docx(path: Path, *, max_chars: int) -> dict[str, Any]:
    try:
        from docx import Document
    except ImportError:
        return _error(
            "python-docx is required. Install with: pip install python-docx",
            max_chars=max_chars,
        )

    try:
        doc = Document(str(path))
    except Exception as e:
        return _error(f"Failed to inspect .docx file: {e}", max_chars=max_chars)

    paragraphs = [str(p.text or "").strip() for p in doc.paragraphs if str(p.text or "").strip()]
    word_count = sum(len(p.split()) for p in paragraphs)
    page_count = max(1, (len(paragraphs) + 29) // 30) if paragraphs else 0

    headings: list[str] = []
    for paragraph in doc.paragraphs:
        text = str(paragraph.text or "").strip()
        style_name = str(getattr(getattr(paragraph, "style", None), "name", "") or "").strip().lower()
        if not text or not style_name.startswith("heading"):
            continue
        match = re.search(r"(\d+)", style_name)
        level = int(match.group(1)) if match else 1
        level = min(6, max(1, level))
        headings.append(f"{'  ' * (level - 1)}- {text}")

    lines = [
        f"Path: {path}",
        "Type: DOCX",
        f"Approximate pages: {page_count}",
        f"Word count: {word_count}",
        "Heading outline:",
    ]
    if headings:
        lines.extend(headings)
    else:
        lines.append("- (none)")
    return _success("\n".join(lines), max_chars=max_chars)


def _info_xlsx(path: Path, *, max_chars: int) -> dict[str, Any]:
    try:
        from openpyxl import load_workbook
    except ImportError:
        return _error(
            "openpyxl is required. Install with: pip install openpyxl",
            max_chars=max_chars,
        )

    try:
        workbook = load_workbook(filename=str(path), data_only=True, read_only=True)
    except Exception as e:
        return _error(f"Failed to inspect .xlsx file: {e}", max_chars=max_chars)

    try:
        lines: list[str] = [f"Path: {path}", "Type: XLSX", "Sheets:"]
        for ws in workbook.worksheets:
            header_values = next(ws.iter_rows(min_row=1, max_row=1, values_only=True), ())
            headers = [str(v) if v is not None else "" for v in header_values]
            lines.append(f"- {ws.title}: rows={int(ws.max_row or 0)}, cols={int(ws.max_column or 0)}")
            if headers:
                lines.append(f"  headers: {', '.join(headers)}")
            else:
                lines.append("  headers: (none)")
        return _success("\n".join(lines), max_chars=max_chars)
    finally:
        try:
            workbook.close()
        except Exception:
            pass


def _info_pptx(path: Path, *, max_chars: int) -> dict[str, Any]:
    try:
        from pptx import Presentation
    except ImportError:
        return _error(
            "python-pptx is required. Install with: pip install python-pptx",
            max_chars=max_chars,
        )

    try:
        presentation = Presentation(str(path))
    except Exception as e:
        return _error(f"Failed to inspect .pptx file: {e}", max_chars=max_chars)

    lines: list[str] = [
        f"Path: {path}",
        "Type: PPTX",
        f"Slide count: {len(presentation.slides)}",
        "Slide titles:",
    ]
    for index, slide in enumerate(presentation.slides, start=1):
        title = ""
        try:
            if slide.shapes.title is not None:
                title = str(slide.shapes.title.text or "").strip()
        except Exception:
            title = ""
        lines.append(f"- Slide {index}: {title or 'Untitled'}")
    return _success("\n".join(lines), max_chars=max_chars)


def _modify_xlsx(path: Path, *, sheet_name: str | None, changes: Any, max_chars: int) -> dict[str, Any]:
    try:
        from openpyxl import load_workbook
    except ImportError:
        return _error(
            "openpyxl is required. Install with: pip install openpyxl",
            max_chars=max_chars,
        )

    if changes is None:
        return _error("changes is required for .xlsx modify action", max_chars=max_chars)

    try:
        updates = json.loads(changes) if isinstance(changes, str) else changes
    except Exception as e:
        return _error(f"changes must be valid JSON: {e}", max_chars=max_chars)

    if not isinstance(updates, list):
        return _error("changes must be a JSON array of {row, col, value}", max_chars=max_chars)

    try:
        workbook = load_workbook(filename=str(path))
    except Exception as e:
        return _error(f"Failed to open .xlsx file: {e}", max_chars=max_chars)

    try:
        if sheet_name:
            if sheet_name not in workbook.sheetnames:
                return _error(
                    f"Sheet not found: {sheet_name}. Available: {', '.join(workbook.sheetnames)}",
                    max_chars=max_chars,
                )
            ws = workbook[sheet_name]
        else:
            ws = workbook.active

        changed_count = 0
        for update in updates:
            if not isinstance(update, dict):
                return _error("each change must be an object with row, col, value", max_chars=max_chars)
            row = update.get("row")
            col = update.get("col")
            if not isinstance(row, int) or not isinstance(col, int) or row < 1 or col < 1:
                return _error("row and col must be 1-indexed positive integers", max_chars=max_chars)
            value = update.get("value")
            cell = ws.cell(row=row, column=col)
            if cell.value != value:
                cell.value = value
                changed_count += 1

        workbook.save(str(path))
        return _success(f"Applied {changed_count} xlsx cell modifications to {path}", max_chars=max_chars)
    finally:
        try:
            workbook.close()
        except Exception:
            pass


def _modify_docx(
    path: Path,
    *,
    find: str | None,
    replace: str | None,
    max_replacements: int | None,
    max_chars: int,
) -> dict[str, Any]:
    try:
        from docx import Document
    except ImportError:
        return _error(
            "python-docx is required. Install with: pip install python-docx",
            max_chars=max_chars,
        )

    find_text = str(find or "")
    if not find_text:
        return _error("find is required for .docx modify action", max_chars=max_chars)
    replace_text = str(replace or "")

    limit: int | None
    if max_replacements is None or int(max_replacements) <= 0:
        limit = None
    else:
        limit = int(max_replacements)

    try:
        doc = Document(str(path))
    except Exception as e:
        return _error(f"Failed to open .docx file: {e}", max_chars=max_chars)

    count = 0
    remaining = limit
    for paragraph in doc.paragraphs:
        text = str(paragraph.text or "")
        if not text or find_text not in text:
            continue

        occurrences = text.count(find_text)
        if occurrences <= 0:
            continue

        if remaining is None:
            paragraph.text = text.replace(find_text, replace_text)
            count += occurrences
        else:
            local_count = min(occurrences, remaining)
            paragraph.text = text.replace(find_text, replace_text, local_count)
            count += local_count
            remaining -= local_count
            if remaining <= 0:
                break

    try:
        if count > 0:
            doc.save(str(path))
    except Exception as e:
        return _error(f"Failed to save .docx file: {e}", max_chars=max_chars)

    return _success(f"Applied {count} replacements to {path}", max_chars=max_chars)


@tool
def office(
    action: str = "read",
    path: str = "",
    *,
    content: Any = None,
    sheet_name: str | None = None,
    max_rows: int = 200,
    max_chars: int = 12000,
    changes: Any = None,
    find: str | None = None,
    replace: str | None = None,
    max_replacements: int | None = None,
    cwd: str | None = None,
) -> dict[str, Any]:
    """
    Read and write Office files (.docx/.xlsx/.pptx) using format-aware extraction and generation.
    """
    mode = (action or "read").strip().lower()
    if mode not in {"read", "write", "info", "modify"}:
        return _error(f"Unknown action: {mode}", max_chars=max_chars)

    for_write = mode == "write"
    try:
        _, target = _resolve_path(path, cwd=cwd, for_write=for_write)
    except Exception as e:
        return _error(str(e), max_chars=max_chars)

    ext = target.suffix.lower()
    if ext not in _SUPPORTED_EXTENSIONS:
        return _error(
            f"Unsupported Office format: {ext or '(none)'}. Supported: .docx, .xlsx, .pptx",
            max_chars=max_chars,
        )

    if mode == "read":
        if ext == ".docx":
            return _read_docx(target, max_chars=max_chars)
        if ext == ".xlsx":
            return _read_xlsx(target, sheet_name=sheet_name, max_rows=max_rows, max_chars=max_chars)
        return _read_pptx(target, max_chars=max_chars)

    if mode == "write":
        if ext == ".docx":
            return _write_docx(target, content=content, max_chars=max_chars)
        if ext == ".xlsx":
            return _write_xlsx(target, content=content, sheet_name=sheet_name, max_chars=max_chars)
        return _write_pptx(target, content=content, max_chars=max_chars)

    if mode == "info":
        if ext == ".docx":
            return _info_docx(target, max_chars=max_chars)
        if ext == ".xlsx":
            return _info_xlsx(target, max_chars=max_chars)
        return _info_pptx(target, max_chars=max_chars)

    if ext == ".xlsx":
        return _modify_xlsx(target, sheet_name=sheet_name, changes=changes, max_chars=max_chars)
    return _modify_docx(
        target,
        find=find,
        replace=replace,
        max_replacements=max_replacements,
        max_chars=max_chars,
    )


set_permissions(office, "read")
